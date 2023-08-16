from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_VERTICAL_ANCHOR as MSO_ANCHOR
from pptx.dml.color import RGBColor as RGB

def create_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    subtitle_placeholder.text = subtitle

def add_model_slide(prs, model_name, pros, cons, font_size=18):
    slide_layout = prs.slide_layouts[5]  # Blank slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title_shape = slide.shapes.title
    title_shape.text = model_name

    # Define dimensions for the text box
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(4)

    # Create a text box shape with the specified dimensions
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP

    # Adding pros
    pros_title = tf.add_paragraph()
    pros_title.text = "Pros"
    pros_title.font.bold = True
    pros_title.font.size = Pt(20)  # Font size for the "Pros" title
    pros_title.font.color.rgb = RGB(0, 128, 0)

    for pro in pros:
        p = tf.add_paragraph()
        p.text = "- " + pro
        p.level = 1
        for run in p.runs:
            run.font.size = Pt(font_size)  # Set the font size

    # Adding spacing
    tf.add_paragraph()

    # Adding cons
    cons_title = tf.add_paragraph()
    cons_title.text = "Cons"
    cons_title.font.bold = True
    cons_title.font.size = Pt(20)
    cons_title.font.color.rgb = RGB(255, 0, 0)

    for con in cons:
        p = tf.add_paragraph()
        p.text = "- " + con
        p.level = 1
        for run in p.runs:
            run.font.size = Pt(font_size)  # Set the font size
def generate_presentation(filename):
    prs = Presentation()

    # Title Slide
    create_title_slide(prs, "Model Comparison for Ethane Production Forecasting", "Analyzing the Pros and Cons based on the Pipeline Scenario")

    # LSTM Model

    # ARIMA Model
    add_model_slide(
        prs,
        "ARIMA",
        ["Handles trend and seasonality: Good for capturing patterns in ethane production",
         "Widely used and understood: Familiarity within forecasting communities",
         "Works well with stationary data: Can work with transformed ethane production data"],
        ["Requires stationary data: Ethane data may need differencing or transformations",
         "Manual parameter tuning can be complex: Needs expertise for ethane-specific variations",
         "Struggles with high-frequency data: Daily fluctuations in ethane data can be challenging"],
        font_size=14
    )

    # Prophet Model
    add_model_slide(
        prs,
        "Prophet",
        ["Handles holidays and special events: Can consider fluctuations during maintenance or disruptions",
         "Robust to missing data: Tolerant to occasional missing records in ethane data",
         "Automatic parameter tuning: Less manual intervention needed"],
        ["Requires daily data: Might need consistent data recording for ethane production",
         "Less interpretable parameters: Black-box nature might not provide insights on ethane production specifics",
         "May overfit if seasonality is too strong: Ethane production's strong seasonal patterns need careful handling"],
        font_size=14
    )

    # SARIMAX Model
    add_model_slide(
        prs,
        "SARIMAX",
        ["Handles trend, seasonality, and exogenous variables: Can integrate price spreads, weather data, etc. for ethane forecasting",
         "Can model multiple seasonalities: Daily and yearly fluctuations in ethane production",
         "Good for long time series: Suitable for extensive ethane production records"],
        ["Requires stationary data: Ethane data might need transformations",
         "Complex parameter tuning: Need expertise for ethane-specific intricacies",
         "Computationally expensive: Might be slow for large ethane datasets"],
        font_size=14
    )
    add_model_slide(
        prs,
        "LSTM",
        ["Memory of Past Data: Captures long-term dependencies", "Adaptable to Non-linearity: Models complex patterns", "Handles Multiple Variables: Integrates price, weather, etc."],
        ["Requires Large Data: Needs extensive data for optimal performance", "Training Time: Computationally intensive", "Complexity: Parameter tuning can be challenging"],
        font_size=14
    )

    # XGBoost Model
    add_model_slide(
        prs,
        "XGBoost",
        ["Handles Non-linearity: Captures diverse patterns", "Feature Importance: Insights into influential variables", "Robustness: Less sensitive to outliers"],
        ["Requires Feature Engineering: Needs careful variable crafting", "Overfitting: Can over-adapt to training data without tuning"],
        font_size=14
    )

    # VAR Model
    add_model_slide(
        prs,
        "Vector Autoregression (VAR)",
        ["Multiple Time Series: Models interdependencies of variables", "Interpretability: Insights from model coefficients"],
        ["Stationarity Requirement: Needs all time series to be stationary", "Model Complexity: Can become intricate with many variables"],
        font_size=14
    )

    # TBATS Model
    add_model_slide(
        prs,
        "TBATS",
        ["Multiple Seasonalities: Handles daily and yearly patterns", "Automatic Box-Cox Transformation: Stabilizes variance"],
        ["Computational Intensity: Can be resource-consuming", "Black-box Nature: Less transparent than some models"],
        font_size=14
    )

    prs.save(filename)

generate_presentation("models_comparison_presentation.pptx")
