from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor as RGB

# Load presentation
prs = Presentation("models_comparison_presentation.pptx")

# Define colors
dark_blue = RGB(34, 66, 124)
light_blue = RGB(91, 155, 213)

# Function to set font styles for a shape
def style_text_frame(tf, font_size=18, font_color=dark_blue):
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Arial'
            run.font.size = Pt(font_size)
            run.font.color.rgb = font_color

# Function to add bullet points
def add_bullets(tf):
    tf.text = ''  # clear any existing text
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = '• Point 1'
    p.level = 1
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = '• Point 2'
    p.level = 1
    style_text_frame(tf)

# Iterate through all slides and apply styles
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        tf.margin_left = Inches(0.5)
        tf.margin_right = Inches(0.5)

        # Check for "Pros" and "Cons" shapes and add bullet points
        title_text = shape.text_frame.text
        if title_text == "Pros" or title_text == "Cons":
            add_bullets(tf)
            style_text_frame(tf, font_size=24, font_color=light_blue)
            shape.fill.solid()
            shape.fill.fore_color.rgb = dark_blue
            shape.line.width = Inches(0.05)
            shape.line.color.rgb = dark_blue
        else:
            style_text_frame(tf)

# Save the presentation
prs.save("styled_models_comparison_presentation.pptx")
