from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_VERTICAL_ANCHOR as MSO_ANCHOR
from pptx.dml.color import RGBColor as RGB


# Create title slide
def create_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    subtitle_placeholder.text = subtitle
# Add model slide with detailed explanations
# Add combined model explanations as a table
def add_combined_explanation_slide(prs, models):
    slide_layout = prs.slide_layouts[5]  # Blank slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title_shape = slide.shapes.title
    title_shape.text = "Model Explanations"

    # Add a table to the slide
    rows, cols = len(models) + 1, 2  # Including header row
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column names
    table.cell(0, 0).text = "Model"
    table.cell(0, 1).text = "Explanation"

    # Fill the table with model data
    for idx, model in enumerate(models, start=1):
        model_name, _, _, explanation = model
        table.cell(idx, 0).text = model_name
        table.cell(idx, 1).text = explanation
    # Adjust font size and cell margins for better fitting
    for row in table.rows:
        row.height = Inches(0.5)
        for cell in row.cells:
            cell.text_frame.margin_top = Inches(0.1)
            cell.text_frame.margin_bottom = Inches(0.1)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)



# Add model slide with detailed pros and cons and explanations
def add_model_detail_slide(prs, model_name, explanation):
    slide_layout = prs.slide_layouts[5]  # Blank slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title_shape = slide.shapes.title
    title_shape.text = model_name

    # Model explanation
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(3)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.add_paragraph()
    p.text = explanation
    p.font.size = Pt(14)  # Adjusted font size

# Create table slide with concise pros and cons
def create_summary_table_slide(prs, models):
    slide_layout = prs.slide_layouts[5]  # Blank slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title_shape = slide.shapes.title
    title_shape.text = "Models Overview"

    # Add a table to the slide
    rows, cols = len(models) + 1, 3  # Including header row
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column names
    table.cell(0, 0).text = "Model"
    table.cell(0, 1).text = "Pros"
    table.cell(0, 2).text = "Cons"

    # Fill the table with model data
    for idx, model in enumerate(models, start=1):
        model_name, pros, cons, _ = model
        table.cell(idx, 0).text = model_name
        table.cell(idx, 1).text = '\n'.join(["• " + pro for pro in pros])  # Use bullet points
        table.cell(idx, 2).text = '\n'.join(["• " + con for con in cons])  # Use bullet points

    # Adjust font size and cell margins for better fitting
    for row in table.rows:
        row.height = Inches(0.5)
        for cell in row.cells:
            cell.text_frame.margin_top = Inches(0.1)
            cell.text_frame.margin_bottom = Inches(0.1)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)  # Set font size to 12 for all cells



def generate_presentation(filename):
    prs = Presentation()

    # Title Slide
    create_title_slide(prs, "Model Comparison for Ethane Production Forecasting", "Analyzing the Pros and Cons based on the Pipeline Scenario")

    # Model data: (model_name, pros, cons, explanation)
    models_data = [
        ("ARIMA",
         ["Works well with data that has clear patterns", "Best for stable data"],
         ["Needs stable data", "Requires manual setup", "Not for frequently changing data"],
         "Think of ARIMA as a classic recipe for predicting patterns in data. It works best when the data doesn't "
         "have too many unexpected ups and downs."
         ),

        ("Prophet",
         ["Accounts for special occasions like holidays", "Works with missing data", "Automatic setup"],
         ["Best with daily data", "Black-box nature", "Can over-predict details"],
         "Prophet, made by Facebook, is like a weather forecast for data. It predicts daily patterns and even "
         "considers holidays."
         ),

        ("SARIMAX",
         ["Predicts based on multiple patterns and external info", "Understands multiple seasons", "Great for longer datasets"],
         ["Needs stable data", "Can be tricky to set up", "Computationally demanding"],
         "Think of SARIMAX as an upgraded version of ARIMA. It can use extra information and understand more "
         "complicated patterns."
         ),

        ("LSTM",
         ["Remembers past data for a long time", "Adapts to unexpected changes", "Works with multiple data types"],
         ["Needs lots of past data", "Takes time to set up", "Complex to fine-tune"],
         "LSTM is like a diary that remembers important events from the past to predict the future."
         ),

        ("XGBoost",
         ["Handles unexpected changes", "Identifies important info", "Not fooled by outliers"],
         ["Needs data setup", "Can over-predict details", "Requires tuning"],
         "XGBoost is like a high-tech detective. It's known to find patterns better than most other methods."
         ),

        ("VAR",
         ["Predicts using multiple datasets", "Insights on mutual data influence"],
         ["Needs all stable data", "Complex with too much info"],
         "VAR observes multiple datasets, like understanding how one dataset's behavior affects the others."
         ),

        ("TBATS",
         ["Understands multiple seasonal patterns", "Automatic setup"],
         ["Computationally intensive", "Black-box nature"],
         "TBATS is designed to understand data with recurring patterns, like a song with repeating choruses."
         )
    ]



    # Create summary table slide
    create_summary_table_slide(prs, models_data)

    # Create a single slide with model explanations
    add_combined_explanation_slide(prs, models_data)

    prs.save(filename)

generate_presentation("models_comparison_presentationv2.pptx")
